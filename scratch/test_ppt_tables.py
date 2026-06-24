from pptx import Presentation

def inspect_shape(shape, idx, indent="  "):
    pos = f"L:{shape.left.inches:.2f} T:{shape.top.inches:.2f} W:{shape.width.inches:.2f} H:{shape.height.inches:.2f}" if hasattr(shape, 'left') else "N/A"
    if shape.has_table:
        print(f"{indent}Shape: {shape.name} [TABLE] ({pos})")
        table = shape.table
        for r_idx, row in enumerate(table.rows):
            row_text = []
            for cell in row.cells:
                row_text.append(cell.text.strip().encode('ascii', 'replace').decode('ascii'))
            print(f"{indent}    Row {r_idx}: {row_text}")
    elif shape.has_text_frame:
        text = shape.text_frame.text.strip()
        print(f"{indent}Shape: {shape.name} [TEXT] ({pos}): {text[:60].encode('ascii', 'replace').decode('ascii')}...")
    elif shape.shape_type == 6: # Group shape
        print(f"{indent}Shape: {shape.name} [GROUP] ({pos})")
        for child in shape.shapes:
            inspect_shape(child, idx, indent + "    ")
    else:
        print(f"{indent}Shape: {shape.name} [OTHER] ({pos}) (Type: {shape.shape_type})")

prs = Presentation("Smart_Verse_Lip2Speech_Presentation.pptx")
for idx, slide in enumerate(prs.slides, 1):
    title = "(No Title Shape)"
    if slide.shapes.title:
        title = slide.shapes.title.text
    else:
        # Search for TextBox shapes containing large text
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                first_line = shape.text_frame.text.split("\n")[0].strip()
                if len(first_line) > 2 and len(first_line) < 50:
                    title = f"(Text search): {first_line}"
                    break
    print(f"Slide {idx:02d}: {title}")
    for shape in slide.shapes:
        inspect_shape(shape, idx)

