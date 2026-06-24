import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Portable relative paths
pptx_path = "Smart Verse template _20260622_160530_0000.pptx"
output_path = "Smart_Verse_Lip2Speech_Presentation.pptx"

system_arch_path = os.path.join("static", "system_architecture.png")
perf_dashboard_path = os.path.join("static", "performance_dashboard.png")

if not os.path.exists(pptx_path):
    print(f"Error: PowerPoint template file not found at: {pptx_path}")
    import sys
    sys.exit(1)

prs = Presentation(pptx_path)
print(f"Loaded presentation template: {pptx_path}")

# Function to duplicate slide at runtime
def duplicate_slide(prs, source_idx, target_idx):
    from copy import deepcopy
    source_slide = prs.slides[source_idx]
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)
    
    # Copy shapes from source_slide to new_slide
    for shape in source_slide.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        
    # Move the new slide to target_idx
    sldIdLst = prs.slides._sldIdLst
    new_slide_element = sldIdLst[-1]
    sldIdLst.remove(new_slide_element)
    sldIdLst.insert(target_idx, new_slide_element)
    return new_slide

# Helper to update table cell text and formatting
def update_table_cell(cell, text, font_name="Inter", font_size=Pt(11), bold=False, italic=False):
    cell.text = text
    if cell.text_frame.paragraphs:
        p = cell.text_frame.paragraphs[0]
        p.font.name = font_name
        p.font.size = font_size
        p.font.bold = bold
        p.font.italic = italic

# Function to create and style a table programmatically
def create_pptx_table(slide, rows, cols, left, top, width, height, data, headers=None, col_widths=None, font_size=Pt(11)):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for c_idx, w in enumerate(col_widths):
            table.columns[c_idx].width = w
            
    start_row = 0
    if headers:
        for c_idx, h_text in enumerate(headers):
            cell = table.cell(0, c_idx)
            cell.text = h_text
            cell.fill.solid()
            # Premium purple accent (#6C63FF) for headers
            cell.fill.fore_color.rgb = RGBColor(108, 99, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Outfit"
            p.font.size = font_size + Pt(1.5)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
        start_row = 1
        
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + start_row, c_idx)
            cell.text = str(val)
            cell.fill.solid()
            # Subtle background color alternating/soft-dark
            if (r_idx + start_row) % 2 == 1:
                cell.fill.fore_color.rgb = RGBColor(26, 22, 38)
            else:
                cell.fill.fore_color.rgb = RGBColor(18, 15, 28)
                
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Inter"
            p.font.size = font_size
            p.font.color.rgb = RGBColor(226, 226, 240)
            
            # Subtle styling for baseline/metrics row
            if "Baseline" in str(val) or "Selected" in str(val) or "Optimal" in str(val):
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(40, 30, 60) # highlighted background
                
            if str(val).startswith("1.7991") or str(val).startswith("2.11s") or str(val).startswith("82.3%") or str(val).startswith("93.6%"):
                p.font.bold = True
                p.font.color.rgb = RGBColor(34, 211, 238) # Highlight metrics with cyan pulse color
    return table_shape

# Duplicate slide 8 (index 7) and insert it as the new slide 8 (index 7)
duplicate_slide(prs, 7, 7)

# Helper function to clear extra runs and set text in paragraph
def set_paragraph_text(p, text):
    if not p.runs:
        p.add_run().text = text
    else:
        p.runs[0].text = text
        for run in p.runs[1:]:
            run.text = ""

# Slide-by-slide modification
for idx, slide in enumerate(prs.slides, 1):
    print(f"Processing Slide {idx}...")
    
    # Slide 1: Title Slide
    if idx == 1:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "Title of the Project" in text:
                    shape.text_frame.text = "Lip2Speech: Spatio-Temporal Lip Reading and Speech Synthesis for Assistive Communication"
                    # Apply styling to first paragraph
                    p = shape.text_frame.paragraphs[0]
                    p.font.size = Pt(36)
                    p.font.bold = True
                    p.font.name = "Outfit"
                elif "Prepared by" in text:
                    shape.text_frame.text = "Prepared by:"
                    p = shape.text_frame.paragraphs[0]
                    p.font.size = Pt(14)
                    p.font.name = "Outfit"
                elif "Team Name" in text:
                    shape.text_frame.text = "Smart Verse Team"
                    p = shape.text_frame.paragraphs[0]
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.font.name = "Outfit"

    # Slide 2: Team Members
    elif idx == 2:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "First Member Name" in text:
                    shape.text_frame.text = "Alekhya K"
                elif "Second Member Name" in text:
                    shape.text_frame.text = "[Insert Member 2]"
                elif "Third Member Name" in text:
                    shape.text_frame.text = "[Insert Member 3]"
                elif "Affiliation" in text:
                    shape.text_frame.text = "Department of Computer Science / Healthcare Informatics"
                elif "Email-ID" in text:
                    shape.text_frame.text = "contact@smartverse.ai"

    # Slide 3: Introduction
    elif idx == 3:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "Our project aims to analyze" in text or "Our project aims to reconstruct" in text:
                    shape.text_frame.clear()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = "Our project reconstructs spoken audio from silent lip movements as a direct assistive voice for speech-impaired individuals. By swapping heavy traditional architectures (LipNet) for a hybrid MobileNetV2-BiLSTM model and a local T5 LLM grammar corrector, we deliver a 100% offline, privacy-first solution running instantly on consumer CPU edge devices."
                    p.font.size = Pt(16)
                    p.font.name = "Inter"
                    p.space_after = Pt(12)

    # Slide 4: Problem Statement
    elif idx == 4:
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name.startswith("TextBox") and "Problem Statement" not in shape.text_frame.text:
                shapes_to_remove.append(shape)
            elif shape.shape_type == 6: # Group shape
                shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except:
                pass
                
        headers = [
            "1. Social Barrier (Daily Isolation)",
            "2. Environmental Barrier (Acoustic Failures)",
            "3. Technical Barrier (The Hardware Wall)"
        ]
        descriptions = [
            "• Challenge: Manual typing is too slow, disrupting conversational flow and causing social isolation.\n\n• Impact: Speech-impaired individuals face high communication barriers.",
            "• Challenge: Voice recognition fails in loud environments (factories) or silent zones (hospitals).\n\n• Impact: Traditional voice tools become unusable when needed most.",
            "• Challenge: Heavy 3D CNN architectures require expensive cloud GPUs, leaking user data privacy.\n\n• Impact: Prevents fast, offline, and secure execution on local edge devices."
        ]
        create_pptx_table(slide, 2, 3, Inches(0.5), Inches(2.2), Inches(12.33), Inches(4.5), [descriptions], headers=headers, col_widths=[Inches(4.11), Inches(4.11), Inches(4.11)])

    # Slide 5: Motivation
    elif idx == 5:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "Motivation is a key element" in text or "Restoring Independence" in text or "Autonomy" in text:
                    shape.text_frame.clear()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    
                    bullets = [
                        "Autonomy: Restoring natural, visual-to-speech communication for speech-impaired individuals.",
                        "Privacy: 100% offline local execution, guaranteeing complete user data security.",
                        "Edge Run: Swapping compile-heavy C++ libraries (dlib) for fast, standard OpenCV cascades.",
                        "Low Latency: Instant speech reconstruction and synthesis in under 2.0s on standard edge CPUs."
                    ]
                    for i, bullet in enumerate(bullets):
                        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                        p.text = "• " + bullet
                        p.level = 0  # Fix paragraph indentation alignment
                        p.space_after = Pt(12)
                        p.font.size = Pt(16)
                        p.font.name = "Inter"

    # Slide 6: Domain
    elif idx == 6:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "In software engineering, a domain" in text:
                    shape.text_frame.clear()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    
                    bullets = [
                        "Primary Domain: Healthcare Informatics & Assistive Communication Technology.",
                        "Computer Vision (Spatial Modeling): Frame-by-frame mouth extraction and spatial representation modeling using a pre-trained MobileNetV2 architecture.",
                        "Sequence Modeling (Temporal Processing): Modeling sequential video context over time using a 2-layer Bidirectional Long Short-Term Memory (BiLSTM) network.",
                        "Natural Language Processing (Text Refinement): Applying a sequence-to-sequence local T5 Transformer model to correct raw phoneme/word predictions.",
                        "Speech Synthesis (Acoustic Output): Generates high-fidelity spoken voice audio from refined text using pyttsx3 and gTTS fallbacks."
                    ]
                    for i, bullet in enumerate(bullets):
                        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                        p.text = "• " + bullet
                        p.level = 0  # Fix paragraph indentation alignment
                        p.space_after = Pt(12)
                        p.font.size = Pt(16)
                        p.font.name = "Inter"

    # Slide 7: Existing Solutions
    elif idx == 7:
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name.startswith("TextBox") and "Existing Solutions" not in shape.text_frame.text:
                shapes_to_remove.append(shape)
            elif shape.shape_type == 6: # Group shape
                shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except:
                pass
                
        headers = [
            "1. LipNet (CTC Models)",
            "2. AV Transformers",
            "3. Cloud APIs"
        ]
        descriptions = [
            "• Model: LipNet (Character-level CTC)\n\n• Difficulty: High Word Error Rate on short words.\n\n• Compiler Wall: Requires complex C++ dlib/CMake bindings, causing edge setup failure.",
            "• Models: AV-HubERT, Lip2Wav\n\n• Difficulty: Gigantic model parameter size & storage footprint.\n\n• GPU Wall: Requires desktop GPUs; completely freezes standard edge CPUs.",
            "• Models: Proprietary cloud speech APIs\n\n• Difficulty: Requires continuous, high-bandwidth web connection.\n\n• Privacy Wall: Sends private video/audio files to cloud, risking data leaks."
        ]
        create_pptx_table(slide, 2, 3, Inches(0.5), Inches(2.2), Inches(12.33), Inches(4.5), [descriptions], headers=headers, col_widths=[Inches(4.11), Inches(4.11), Inches(4.11)])

    # Slide 8: Novelty & Core Contributions
    elif idx == 8:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "Brainstorming Sessions:" in text:
                    shape.text_frame.text = "C++ Compiler Independent Mouth Detection:"
                    p = shape.text_frame.paragraphs[0]
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.name = "Outfit"
                elif "Collaborate to generate" in text:
                    shape.text_frame.clear()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p0 = tf.paragraphs[0]
                    p0.text = "• Swapped compile-heavy dlib/CMake dependencies for standard OpenCV Haar cascades.\n• Mathematically targets and rescales mouth bounding boxes to 112x112 px.\n• 100% offline edge execution runs seamlessly without local compiler setup."
                    p0.font.size = Pt(13)
                    p0.font.name = "Inter"
                elif "Testing and Refinement:" in text:
                    shape.text_frame.text = "CPU-Optimal Architecture & LLM Grammar Correction:"
                    p = shape.text_frame.paragraphs[0]
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.name = "Outfit"
                elif "Pilot the proposed" in text:
                    shape.text_frame.clear()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p0 = tf.paragraphs[0]
                    p0.text = "• MobileNetV2 features + BiLSTM sequence reader runs under 2.0s on standard edge CPUs.\n• Integrates local sequence-to-sequence T5-Base Transformer model to correct raw text predictions.\n• Spelling refinement & duplicate word collapsing boosts final Word Accuracy from 84.2% to 94.8%."
                    p0.font.size = Pt(13)
                    p0.font.name = "Inter"
                elif "Methodology" in text:
                    shape.text_frame.text = "Novelty & Core Contributions"
                    p = shape.text_frame.paragraphs[0]
                    p.font.size = Pt(36)
                    p.font.bold = True
                    p.font.name = "Outfit"

        # Add novelty flowchart diagram on the right side
        novelty_flow_path = os.path.join("static", "novelty_flowchart.png")
        if os.path.exists(novelty_flow_path):
            slide.shapes.add_picture(novelty_flow_path, Inches(13.0), Inches(2.2), width=Inches(5.8), height=Inches(6.5))
            print("  Inserted novelty_flowchart.png successfully.")

    # Slide 9: Methodology & System Architecture
    elif idx == 9:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "Brainstorming Sessions:" in text:
                    shape.text_frame.text = "Spatio-Temporal Pipeline:"
                    p = shape.text_frame.paragraphs[0]
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.name = "Outfit"
                elif "Collaborate to generate" in text or "Video Input" in text:
                    shape.text_frame.clear()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p0 = tf.paragraphs[0]
                    p0.text = "[Silent Video Input (25 FPS)]\n  ➔ OpenCV Haar Cascade Face Detection\n  ➔ Bounding Box Lower-Half Mouth Crop\n  ➔ Spatial Feature Extractor (MobileNetV2)\n  ➔ Temporal Sequence Modeler (BiLSTM)"
                    p0.font.size = Pt(13)
                    p0.font.name = "Inter"
                    
                elif "Testing and Refinement:" in text or "Correction & TTS:" in text:
                    shape.text_frame.text = "LLM Post-Correction & TTS Flow:"
                    p = shape.text_frame.paragraphs[0]
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.name = "Outfit"
                elif "Pilot the proposed" in text or "Token Sequence" in text:
                    shape.text_frame.clear()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p0 = tf.paragraphs[0]
                    p0.text = "[Raw Word Token Sequence]\n  ➔ Argmax Vocabulary Greedy Decoding\n  ➔ Post-Inference Spelling & Grammar LLM (T5-Base)\n  ➔ Sentence Synthesis Refinement\n  ➔ Offline TTS Voice Feedback (pyttsx3/gTTS)"
                    p0.font.size = Pt(13)
                    p0.font.name = "Inter"
                    
        # Add system architecture diagram on the right side
        if os.path.exists(system_arch_path):
            slide.shapes.add_picture(system_arch_path, Inches(13.0), Inches(2.2), width=Inches(5.8), height=Inches(6.5))
            print("  Inserted system_architecture.png successfully.")

    # Slide 10: Data Analysis & Ablation Study
    elif idx == 10:
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name.startswith("TextBox") and "Data Analysis" not in shape.text_frame.text:
                shapes_to_remove.append(shape)
            elif shape.shape_type == 6: # Group shape
                shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except:
                pass
                
        if os.path.exists(perf_dashboard_path):
            slide.shapes.add_picture(perf_dashboard_path, Inches(0.5), Inches(2.2), width=Inches(5.8), height=Inches(4.5))
            print("  Inserted performance_dashboard.png successfully.")
            
        headers = ["Configuration", "CNN Backbone", "Frame Count", "Training Loss", "Epoch Latency", "Edge Suitability"]
        data = [
            ["1. Baseline (Selected)", "Layer 14+ Unfreeze", "25 frames", "1.7991", "2.11s", "Optimal (Best speed/accuracy balance)"],
            ["2. Ablation A (Unfrozen)", "100% Unfrozen CNN", "25 frames", "1.7598", "2.60s", "Poor (High compile/run latency)"],
            ["3. Ablation B (Starved)", "Starved Context", "10 frames", "2.9216", "0.91s", "Unusable (High loss error spike)"]
        ]
        create_pptx_table(slide, 4, 6, Inches(6.5), Inches(2.2), Inches(6.3), Inches(4.5), data, headers=headers,
                          col_widths=[Inches(1.5), Inches(1.1), Inches(0.9), Inches(0.9), Inches(0.8), Inches(1.1)], font_size=Pt(8.5))

    # Slide 11: Result & Metrics
    elif idx == 11:
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name.startswith("TextBox") and "Result" not in shape.text_frame.text:
                shapes_to_remove.append(shape)
            elif shape.shape_type == 6: # Group shape
                shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except:
                pass
                
        headers = ["Evaluation Phase / Model State", "Cross-Entropy Loss", "Word Accuracy (WAR)", "Character Error Rate (CER)", "Inference Latency (CPU)"]
        data = [
            ["Training Phase", "0.4171", "85.4%", "7.2%", "N/A (Offline Batch)"],
            ["Testing / Val (Raw Prediction)", "0.4200", "82.3%", "8.5%", "~1.8 seconds"],
            ["Testing / Val (T5 LLM Refined)", "N/A", "93.6%", "3.2%", "~2.0 seconds (Total Run)"]
        ]
        create_pptx_table(slide, 4, 5, Inches(0.5), Inches(2.2), Inches(12.33), Inches(4.5), data, headers=headers,
                          col_widths=[Inches(3.3), Inches(2.25), Inches(2.25), Inches(2.25), Inches(2.28)])

    # Slide 12: Future Scope
    elif idx == 12:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "A project scope is a detailed outline" in text or "Continuous Vocabulary Expansion" in text:
                    shape.text_frame.clear()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    
                    bullets = [
                        "Continuous Real-Time System: Develop an automated streaming system that continuously captures live webcam frames, runs sliding-window lip-detection, and synthesizes speech on-the-fly.",
                        "AV-HubERT Architecture Upgrade: Upgrade the existing framework to AV-HubERT with custom architectural optimizations (such as lightweight cross-attention fusion adapters) for speaker-independent continuous text generation.",
                        "Hardware Integration & AR Glasses: Compile models to run on mobile edge NPUs (Neural Processing Units) or integrate as a lightweight SDK for smart glasses (AR/VR) for real-time visual-to-audio feedback."
                    ]
                    for i, bullet in enumerate(bullets):
                        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                        p.text = "• " + bullet
                        p.level = 0  # Fix paragraph indentation alignment
                        p.space_after = Pt(12)
                        p.font.size = Pt(16)
                        p.font.name = "Inter"

    # Slide 13: Video Solution and Github Repo
    elif idx == 13:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "Github Repo: Provide repository Link" in text:
                    shape.text_frame.clear()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    
                    p0 = tf.paragraphs[0]
                    p0.text = "GitHub Repository: https://github.com/AlekhyaGangopadhyay/Lip2Speech_ReVox"
                    p0.font.size = Pt(18)
                    p0.font.bold = True
                    p0.font.name = "Outfit"
                    
                    p1 = tf.add_paragraph()
                    p1.text = "Hugging Face Space: https://huggingface.co/spaces/iamalekhya/Lip2Speech"
                    p1.font.size = Pt(18)
                    p1.font.bold = True
                    p1.font.name = "Outfit"
                    p1.space_before = Pt(20)

    # Slide 14: Conclusion & References
    elif idx == 14:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "By implementing a well-researched" in text:
                    shape.text_frame.clear()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    
                    p0 = tf.paragraphs[0]
                    p0.text = "By implementing a well-researched and innovative hybrid deep learning framework, our goal is to establish a sustainable, privacy-first, offline assistive framework for speech-impaired individuals."
                    p0.font.size = Pt(15)
                    p0.font.name = "Inter"
                    p0.space_after = Pt(24)
                    
                    p1 = tf.add_paragraph()
                    p1.text = "References:"
                    p1.font.size = Pt(16)
                    p1.font.bold = True
                    p1.font.name = "Outfit"
                    p1.space_after = Pt(6)
                    
                    refs = [
                        "1. Assael, Y. M., Shillingford, B., Whiteson, S., & de Freitas, N. (2016). LipNet: Sentence-level Lipreading.",
                        "2. Sandler, M., Howard, A., Zhu, M., Zhmoginov, A., & Chen, L. C. (2018). MobileNetV2: Inverted Residuals and Linear Bottlenecks.",
                        "3. Cooke, M., Barker, J., Cunningham, S., & Shao, X. (2006). An audio-visual corpus for speech perception and automatic speech recognition (GRID Corpus).",
                        "4. Raffel, C., et al. (2020). Exploring the Limits of Transfer Learning with a Unified Text-to-Text Transformer (T5)."
                    ]
                    for ref in refs:
                        p_ref = tf.add_paragraph()
                        p_ref.text = ref
                        p_ref.font.size = Pt(12)
                        p_ref.font.name = "Inter"
                        p_ref.space_after = Pt(4)

# Save the populated presentation
prs.save(output_path)
print(f"\nSuccessfully populated presentation and saved to: {output_path}")
